import os, re, requests, base64, io, glob, subprocess
import matplotlib.pyplot as plt
import qrcode  # Requirement: pip install qrcode[pil]
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# --- PATH CONFIGURATION ---
CURRENT_DIR = os.path.dirname(os.path.abspath(__file__))               # Dossier talk/
ROOT_DIR = os.path.abspath(os.path.join(CURRENT_DIR, ".."))            # Racine du projet

BENCH_DIR = os.path.join(ROOT_DIR, "benchmarks")                       # ../benchmarks/
IMG_DIR = os.path.join(ROOT_DIR, "img")                                # ../img/

TALK_FILE = os.path.join(CURRENT_DIR, "crg_architecture_talk.md")      # Fichier local
VIDEO_PATH = os.path.join(IMG_DIR, "CRG_vs_ECS_Simulation.mov")
QR_PATH = os.path.join(IMG_DIR, "simulator_qr.png")

# --- TARGET URLs ---
SIMULATOR_URL = "https://ct-74.github.io/CapabilityRoutingGateway/demo/final_simulator/index.html"

# --- CPPCON THEME COLORS ---
DARK_GREY = RGBColor(15, 15, 15)
CPP_BLUE = RGBColor(100, 150, 255)
CODE_BG = RGBColor(30, 30, 30)
WHITE = RGBColor(220, 220, 220)

# --- ASSET GENERATORS & FETCHERS ---

def generate_simulator_qr():
    """Generates a static QR code pointing to the live simulator"""
    print(f"📡 Generating QR Code for: {SIMULATOR_URL}")
    qr = qrcode.QRCode(version=1, box_size=10, border=4)
    qr.add_data(SIMULATOR_URL)
    qr.make(fit=True)
    img = qr.make_image(fill_color="black", back_color="white")
    if not os.path.exists(IMG_DIR): 
        os.makedirs(IMG_DIR)
    img.save(QR_PATH)
    return QR_PATH

def get_latest_bench_img(prefix):
    """Retrieves the most recent benchmark result image from the bench directory"""
    if not os.path.exists(BENCH_DIR): 
        return None
    files = glob.glob(os.path.join(BENCH_DIR, f"{prefix}_*.png"))
    return max(files, key=os.path.getctime) if files else None

def generate_tes_layout():
    """Generates a matplotlib-based diagram of the ModelShell memory layout"""
    fig, ax = plt.subplots(figsize=(8, 2))
    ax.barh(0, 48, color='#3498db', edgecolor='white', label='Data (T)')
    ax.barh(0, 8, left=48, color='#e74c3c', edgecolor='white', label='Padding')
    ax.barh(0, 8, left=56, color='#2ecc71', edgecolor='white', label='Concept*')
    ax.set_xlim(0, 64); ax.set_axis_off()
    fig.patch.set_facecolor((0.06, 0.06, 0.06))
    buf = io.BytesIO()
    plt.savefig(buf, format='png', transparent=True, dpi=150)
    plt.close()
    return buf

def get_mermaid_img(code):
    """Fetches a rendered Mermaid diagram via the mermaid.ink API"""
    try:
        b64 = base64.b64encode(code.encode('utf-8')).decode('utf-8')
        res = requests.get(f"https://mermaid.ink/img/{b64}", timeout=10)
        return io.BytesIO(res.content) if res.status_code == 200 else None
    except Exception: 
        return None

# --- PRESENTATION ENGINE ---

class PPTXGenerator:
    def __init__(self, lang='EN'):
        self.lang = lang
        self.prs = Presentation()
        # Set widescreen aspect ratio (16:9)
        self.prs.slide_width, self.prs.slide_height = Inches(13.333), Inches(7.5)

    def add_video_to_slide(self, slide, video_path, left=7.2, top=1.8, width=5.5, height=4.0):
        """Integrates a native video file into the specified slide"""
        if os.path.exists(video_path):
            slide.shapes.add_movie(
                video_path, Inches(left), Inches(top), Inches(width), Inches(height),
                poster_frame_image=None, mime_type='video/quicktime'
            )
            return True
        return False

    def create_deck(self, slides_data):
        for data in slides_data:
            slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = DARK_GREY
            
            # --- HIDDEN SLIDE LOGIC (Q&A / BACKUP) ---
            if "[BACKUP]" in data['title'] or "[HIDDEN]" in data['title']:
                # Modifying the underlying XML element to hide the slide in Presentation Mode
                slide._element.set('show', '0')

            # Title Formatting
            title = slide.shapes.title
            # Clean up the title tag for display
            display_title = data['title'].replace("[BACKUP]", "").replace("[HIDDEN]", "").strip()
            title.text = display_title
            p = title.text_frame.paragraphs[0]
            p.font.color.rgb = CPP_BLUE
            p.font.size = Pt(36)

            title_upper = display_title.upper()
            
            is_bench_slide = any(k in title_upper for k in ["BENCHMARK", "RESULTS", "STRESS TEST"])

            # --- SPECIAL LAYOUT: PERFORMANCE SHOWCASE (TABLE + VIDEO + QR) ---
            if is_bench_slide:
                latest_path = get_latest_bench_img("crg_benchmark_final")
                if latest_path:
                    slide.shapes.add_picture(latest_path, Inches(0.5), Inches(1.8), width=Inches(4.5))
                
                self.add_video_to_slide(slide, VIDEO_PATH, left=5.2, top=1.8, width=5.0, height=3.8)
                
                if os.path.exists(QR_PATH):
                    slide.shapes.add_picture(QR_PATH, Inches(10.5), Inches(1.8), width=Inches(2.2))
                    tx_qr = slide.shapes.add_textbox(Inches(10.5), Inches(4.1), Inches(2.2), Inches(1.0))
                    p_qr = tx_qr.text_frame.paragraphs[0]
                    p_qr.text = "Scan to challenge\nthese results live"
                    p_qr.font.size, p_qr.font.color.rgb = Pt(12), WHITE
                    p_qr.alignment = PP_ALIGN.CENTER

            # --- STANDARD LAYOUT: CODE (LEFT) + MEDIA (RIGHT) ---
            else:
                # Code Block Insertion
                if data['code']:
                    box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(6.5), Inches(5.5))
                    box.fill.solid()
                    box.fill.fore_color.rgb = CODE_BG
                    tf = box.text_frame
                    tf.text = data['code']
                    for paragraph in tf.paragraphs:
                        paragraph.font.name, paragraph.font.size, paragraph.font.color.rgb = 'Consolas', Pt(13), WHITE

                # Media Dispatcher
                img_stream = None
                if any(k in title_upper for k in ["LIMIT", "SIMULATION"]):
                    self.add_video_to_slide(slide, VIDEO_PATH)
                elif "DECISION MATRIX" in title_upper or "BREAK-EVEN" in title_upper:
                    p = os.path.join(IMG_DIR, "breakeven_curve.png")
                    if os.path.exists(p): 
                        with open(p, 'rb') as f: img_stream = io.BytesIO(f.read())
                elif "THE HYPERCUBE" in title_upper or "TENSOR" in title_upper:
                    p = os.path.join(IMG_DIR, "tensor_routing.png")
                    if os.path.exists(p): 
                        with open(p, 'rb') as f: img_stream = io.BytesIO(f.read())
                elif data['mermaid']:
                    img_stream = get_mermaid_img(data['mermaid'])
                elif "ModelShell" in display_title:
                    img_stream = generate_tes_layout()

                if img_stream:
                    slide.shapes.add_picture(img_stream, Inches(7.2), Inches(1.8), width=Inches(5.5))

            # Speaker Notes
            notes = data['fr_notes'] if self.lang == 'FR' else data['en_notes']
            slide.notes_slide.notes_text_frame.text = notes

        output_name = f"CRG_Architecture_Talk_{self.lang}.pptx"
        self.prs.save(os.path.join(CURRENT_DIR, output_name))

if __name__ == "__main__":
    generate_simulator_qr()
    
    if os.path.exists(BENCH_DIR):
        try: 
            subprocess.run(["python3", "run_all.py"], cwd=BENCH_DIR, check=True)
        except Exception: 
            print("⚠️ Benchmark update skipped (using existing assets).")

    try:
        with open(TALK_FILE, 'r', encoding='utf-8') as f:
            raw_slides = re.split(r'# SLIDE:', f.read())[1:]
        
        parsed_data = []
        for s in raw_slides:
            parsed_data.append({
                'title': s.split('\n')[0].strip(),
                'code': (re.search(r'## Code\n+```cpp\n+(.*?)\n+```', s, re.S) or [None, None])[1],
                'mermaid': (re.search(r'## Mermaid\n```mermaid\n(.*?)\n```', s, re.S) or [None,None])[1],
                'en_notes': (re.search(r'## EN\n(.*?)(?=\n## FR|\n##|$)', s, re.S) or [None,""])[1].strip(),
                'fr_notes': (re.search(r'## FR\n(.*?)(?=\n##|$)', s, re.S) or [None,""])[1].strip()
            })

        for lang in ['FR', 'EN']:
            print(f"🎨 Generating {lang} Deck...")
            PPTXGenerator(lang).create_deck(parsed_data)
        
        print("✅ Generation complete!")
        
    except Exception as e: 
        print(f"❌ Error during generation: {e}")