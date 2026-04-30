import os, re, requests, base64, io, glob, subprocess
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# --- CONFIGURATION DES CHEMINS ---
CURRENT_DIR = os.path.dirname(os.path.abspath(__file__))
BENCH_DIR = os.path.abspath(os.path.join(CURRENT_DIR, "..", "benchmarks"))
TALK_FILE = os.path.join(CURRENT_DIR, "cppcon_talk.md")
# Chemin vers le dossier contenant tes images et la vidéo
IMG_DIR = os.path.abspath(os.path.join(CURRENT_DIR, "img")) 
VIDEO_PATH = os.path.join(IMG_DIR, "CRG_vs_ECS_Simulation.mov")

# --- COULEURS CPPCON ---
DARK_GREY = RGBColor(15, 15, 15)
CPP_BLUE = RGBColor(100, 150, 255)
CODE_BG = RGBColor(30, 30, 30)
WHITE = RGBColor(220, 220, 220)

# --- GÉNÉRATEURS / FETCHERS ---

def get_latest_bench_img(prefix):
    """Trouve l'image de benchmark la plus récente avec le préfixe donné"""
    if not os.path.exists(BENCH_DIR):
        return None
    files = glob.glob(os.path.join(BENCH_DIR, f"{prefix}_*.png"))
    return max(files, key=os.path.getctime) if files else None

def generate_tes_layout():
    """Génère le schéma du layout mémoire du ModelShell"""
    fig, ax = plt.subplots(figsize=(8, 2))
    ax.barh(0, 48, color='#3498db', edgecolor='white', label='Data (T)')
    ax.barh(0, 8, left=48, color='#e74c3c', edgecolor='white', label='Padding')
    ax.barh(0, 8, left=56, color='#2ecc71', edgecolor='white', label='Concept*')
    ax.set_xlim(0, 64); ax.set_axis_off()
    fig.patch.set_facecolor((0.06, 0.06, 0.06))
    buf = io.BytesIO()
    plt.savefig(buf, format='png', transparent=True, dpi=150)
    plt.close()
    return buf

def get_mermaid_img(code):
    """Rendu Mermaid.ink avec gestion d'erreur"""
    try:
        b64 = base64.b64encode(code.encode('utf-8')).decode('utf-8')
        res = requests.get(f"https://mermaid.ink/img/{b64}", timeout=10)
        return io.BytesIO(res.content) if res.status_code == 200 else None
    except: 
        return None

# --- MOTEUR DE GÉNÉRATION ---

class PPTXGenerator:
    def __init__(self, lang='EN'):
        self.lang = lang
        self.prs = Presentation()
        self.prs.slide_width, self.prs.slide_height = Inches(13.333), Inches(7.5)

    def add_video_to_slide(self, slide, video_path):
        """Insère une vidéo native dans la slide PowerPoint"""
        if os.path.exists(video_path):
            # Positionnement à droite (miroir du bloc de code)
            left, top = Inches(7.2), Inches(1.8)
            width, height = Inches(5.5), Inches(4.0)
            
            # Note: mime_type est important pour les fichiers .mov sur Windows/Office
            slide.shapes.add_movie(
                video_path, 
                left, top, width, height,
                poster_frame_image=None, 
                mime_type='video/quicktime'
            )
            return True
        return False

    def create_deck(self, slides_data):
        for data in slides_data:
            slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = DARK_GREY
            
            # Titre
            title = slide.shapes.title
            title.text = data['title']
            p = title.text_frame.paragraphs[0]
            p.font.color.rgb = CPP_BLUE
            p.font.size = Pt(36)

            # Code (Gauche)
            if data['code']:
                left, top = Inches(0.5), Inches(1.5)
                width, height = Inches(6.5), Inches(5.5)
                
                box = slide.shapes.add_textbox(left, top, width, height)
                box.fill.solid()
                box.fill.fore_color.rgb = CODE_BG
                
                tf = box.text_frame
                tf.text = data['code']
                tf.word_wrap = True
                
                for paragraph in tf.paragraphs:
                    paragraph.font.name = 'Consolas'
                    paragraph.font.size = Pt(13)
                    paragraph.font.color.rgb = WHITE
                    paragraph.alignment = PP_ALIGN.LEFT

            # Média (Droite)
            img_stream = None
            video_inserted = False
            
            # PRIORITÉ 0 : Vidéo pour les slides de simulation/stress test
            title_upper = data['title'].upper()
            if any(k in title_upper for k in ["STRESS TEST", "SIMULATION", "LIMIT"]):
                video_inserted = self.add_video_to_slide(slide, VIDEO_PATH)

            if not video_inserted:
                # PRIORITÉ 1 : Images statiques haute définition
                if "THE HYPERCUBE" in title_upper or "TENSOR" in title_upper:
                    tensor_path = os.path.join(IMG_DIR, "tensor_routing.png")
                    if os.path.exists(tensor_path):
                        with open(tensor_path, 'rb') as f: img_stream = io.BytesIO(f.read())
                elif "DECISION MATRIX" in title_upper or "THRESHOLD" in title_upper:
                    breakeven_path = os.path.join(IMG_DIR, "breakeven_curve.png")
                    if os.path.exists(breakeven_path):
                        with open(breakeven_path, 'rb') as f: img_stream = io.BytesIO(f.read())
                
                # PRIORITÉ 2 : Rendu Mermaid
                elif data['mermaid']:
                    img_stream = get_mermaid_img(data['mermaid'])
                    
                # PRIORITÉ 3 : Schémas automatiques & Benchmarks
                elif "ModelShell" in data['title']:
                    img_stream = generate_tes_layout()
                elif "Benchmark" in data['title']:
                    latest_path = get_latest_bench_img("crg_benchmark_final")
                    if latest_path:
                        with open(latest_path, 'rb') as f: img_stream = io.BytesIO(f.read())

                # Insertion finale de l'image si trouvée
                if img_stream:
                    slide.shapes.add_picture(img_stream, Inches(7.2), Inches(1.8), width=Inches(5.5))

            # Notes de présentation
            notes = data['fr_notes'] if self.lang == 'FR' else data['en_notes']
            slide.notes_slide.notes_text_frame.text = notes

        output_path = os.path.join(CURRENT_DIR, f"CRG_Talk_{self.lang}.pptx")
        self.prs.save(output_path)

# --- MAIN ---

if __name__ == "__main__":
    # 1. Update Benchmarks
    print("🚀 Updating Benchmarks...")
    if os.path.exists(BENCH_DIR):
        try:
            subprocess.run(["python3", "run_all.py"], cwd=BENCH_DIR, check=True)
        except Exception as e:
            print(f"⚠️ Could not update benchmarks: {e}")

    # 2. Parse Talk
    print("📄 Parsing Talk File...")
    try:
        with open(TALK_FILE, 'r', encoding='utf-8') as f:
            raw_slides = re.split(r'# SLIDE:', f.read())[1:]
        
        parsed_data = []
        for s in raw_slides:
            parsed_data.append({
                'title': s.split('\n')[0].strip(),
                'code': (re.search(r'## Code\n+```cpp\n+(.*?)\n+```', s, re.S) or [None, None])[1],
                'mermaid': (re.search(r'## Mermaid\n```mermaid\n(.*?)\n```', s, re.S) or [None,None])[1],
                'en_notes': (re.search(r'## EN\n(.*?)(?=\n##|$)', s, re.S) or [None,""])[1].strip(),
                'fr_notes': (re.search(r'## FR\n(.*?)(?=\n##|$)', s, re.S) or [None,""])[1].strip()
            })

        # 3. Generate Decks
        for lang in ['FR', 'EN']:
            print(f"🎨 Generating {lang} Deck...")
            PPTXGenerator(lang).create_deck(parsed_data)
        
        print("✅ Decks Generated Successfully!")
        
    except Exception as e:
        print(f"❌ Error during generation: {e}")